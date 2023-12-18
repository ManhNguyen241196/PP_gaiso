from Func.Find_slide_menu import classify_slide
from Func.change_color_text_tbl import change_color
from Func.getAll_textbox import get_all_textBox
from Func.render_all_text_menu import render_text_menu

from pptx import Presentation


link_pp = 'C://Users//ADMIN//Desktop//code//python//Tkinter//fornppgaisho//DummyPPGaisho//DummyData.pptx'
path_save = 'C://Users//ADMIN//Desktop//code//python//Tkinter//fornppgaisho//DummyPPGaisho//result.pptx'
    
psr = Presentation(link_pp)
listSlides = psr.slides

#get all data textbox
data_textBox = get_all_textBox(listSlides, {})
print("data_textBox: ", data_textBox) 

#classify slides
(slide_menu_id, slides_part_id) = classify_slide(data_textBox)
print("slide_menu_id: ", slide_menu_id) 
print("slides_part_id: ", slides_part_id) 


# với loại baeng chia 1.5j và 3.0j
    #get table trong  slide menu
obj_menu = render_text_menu(listSlides,slide_menu_id)
print(obj_menu)

    #work in slide part
part_slide = listSlides.get(257)


        #get table in part slide
name_part = ""
list_shap_textBox = []
list_namePart_textBox = []  
list_タ_textBox = []

for shape in part_slide.shapes:
    if shape.has_table:
        tbl = shape.table
        rows = tbl.rows
        check_index = (1,3,5,6,4,7)
        define_row_slide_menu = tbl.cell(1,0).text_frame.text.replace(" ", "").split("-")[0]
        
        object_row_menu = obj_menu[define_row_slide_menu]
        
        #so sánh table part với table menu
        for i in range(1,len(rows)):
            index_list_compare = 1
            for j in check_index:
                #xét từng cell trong bảng
                cell = rows[i].cells[j]
                
                
                cellCheck =rows[i].cells[2]
                #dinh nghia func so sánh
                def compare_func(list_compare):
                        if(cell.is_spanned == True):
                            cell_special = rows[i-1].cells[j]
                            if(cell_special.text_frame.text.replace(" ", "") != list_compare[index_list_compare]):
                                change_color(cell)
                        if(cell.is_spanned == False):
                            if(cell.text_frame.text.replace(" ", "") != list_compare[index_list_compare]):
                                change_color(cell)
                                
                    #check xem dang la 1.5j hay 3.0j               
                if("1.5" in cellCheck.text_frame.text):
                    listCompare = object_row_menu['1.5タ']
                    name_part = listCompare[1]
                    compare_func(listCompare)
                if("3.0" in cellCheck.text_frame.text):
                    listCompare = object_row_menu['3.0タ']
                    compare_func(listCompare)
                index_list_compare =index_list_compare + 1

        #phan loai cac text box
    if shape.has_text_frame:
        content_shape = shape.text_frame.text.replace(" ", "")
        if('#' in content_shape):
            list_shap_textBox.append(shape)
        if((content_shape == "1.5タ") or(content_shape == "3.0タ") ):
           list_タ_textBox.append(shape)
           
        list_texts = ('以下は細部の応力値です', "応力表", 'Table', 'Fig.')
        for text in list_texts:
            if(text in content_shape):
                list_namePart_textBox.append(shape)
                break
            else:
                pass
    print(name_part)

print("list_shap_textBox: ",list_shap_textBox) 
print ("list_namePart_textBox: ", list_namePart_textBox )
print('list_タ_textBox: ', list_タ_textBox)

psr.save(path_save)